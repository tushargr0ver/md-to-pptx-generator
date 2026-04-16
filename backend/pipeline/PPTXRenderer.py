import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from pipeline.LayoutManager import LayoutManager
from pipeline.StorytellerAgent import PresentationStructure, Slide

class PPTXRenderer:
    """
    Renders slides that visually match the reference sample PPTXs.
    Reference patterns observed:
    - 13.33" x 7.5" widescreen
    - Multi-column grid layouts with vertical divider lines
    - Numbered section badges (01, 02, 03...)
    - Horizontal separator line at slide bottom
    - Slide numbers at bottom-right
    - High shape density (17-59 shapes per slide)
    - Brand accent colors for headers and badges
    """
    
    def __init__(self, master_path: str):
        self.layout_manager = LayoutManager(master_path)
        self.prs = self.layout_manager.prs
        
        # Slide dimensions (13.33" x 7.5" widescreen)
        self.SLIDE_W = Inches(13.33)
        self.SLIDE_H = Inches(7.5)
        
        # Brand colors extracted from reference template
        self.COLOR_TITLE = RGBColor(0x1A, 0x1A, 0x2E)      # Dark navy
        self.COLOR_BODY = RGBColor(0x33, 0x33, 0x33)        # Dark gray text
        self.COLOR_SUBTITLE = RGBColor(0x66, 0x66, 0x66)    # Medium gray
        self.COLOR_ACCENT = RGBColor(0xA1, 0x00, 0xFF)      # Purple accent
        self.COLOR_ACCENT2 = RGBColor(0x00, 0xB4, 0xD8)     # Teal accent
        self.COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        self.COLOR_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)    # Light background
        self.COLOR_LINE = RGBColor(0xCC, 0xCC, 0xCC)         # Separator lines
        self.COLOR_BADGE_BG = RGBColor(0xA1, 0x00, 0xFF)    # Badge background
        
        # Chart color palette
        self.CHART_COLORS = [
            RGBColor(0xA1, 0x00, 0xFF),  # Purple
            RGBColor(0x00, 0xB4, 0xD8),  # Teal
            RGBColor(0xFF, 0x6B, 0x6B),  # Coral
            RGBColor(0x51, 0xCF, 0x66),  # Green
            RGBColor(0xFF, 0xD4, 0x3B),  # Yellow
            RGBColor(0x84, 0x5E, 0xF7),  # Light purple
        ]
        
        # Margins (matching reference)
        self.MARGIN_LEFT = Inches(0.38)
        self.MARGIN_RIGHT = Inches(0.38)
        self.CONTENT_TOP = Inches(1.55)
        self.CONTENT_WIDTH = Inches(12.57)  # 13.33 - 0.38*2
        
        self.slide_number = 0
        
    def render_slides(self, presentation_data: PresentationStructure) -> io.BytesIO:
        for slide_data in presentation_data.slides:
            self.slide_number += 1
            self._render_single_slide(slide_data)
            
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

    def _render_single_slide(self, slide_data: Slide):
        layout = self.layout_manager.get_layout(slide_data.slide_type)
        slide = self.prs.slides.add_slide(layout)
        
        # Remove placeholders we are not using to prevent prompt text (like "Source") from appearing
        unused_ph_indices = []
        is_title_slide = (slide_data.slide_type == "title_slide")
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            ptype = ph.placeholder_format.type
            
            # Keep Title (1) and Slide Number (13) everywhere.
            # Keep indices 10 and 11 ONLY on the Title Slide.
            keep = (ptype in [1, 13]) or (is_title_slide and idx in [10, 11])
            
            if not keep:
                unused_ph_indices.append(idx)
        
        # We must iterate over a list of shapes to delete them safely
        for shape in list(slide.shapes):
            if shape.is_placeholder and shape.placeholder_format.idx in unused_ph_indices:
                sp = shape._element
                sp.getparent().remove(sp)
        
        if slide_data.slide_type == "title_slide":
            self._render_title_slide(slide, slide_data)
        elif slide_data.slide_type in ["content_text", "bullet_points"]:
            self._render_content_slide(slide, slide_data)
        elif slide_data.slide_type == "content_chart":
            self._render_chart_slide(slide, slide_data)
        elif slide_data.slide_type == "infographic_process":
            self._render_infographic_slide(slide, slide_data)
        elif slide_data.slide_type == "infographic_swot":
            self._render_swot_slide(slide, slide_data)
        elif slide_data.slide_type == "infographic_comparison":
            self._render_comparison_slide(slide, slide_data)
        elif slide_data.slide_type == "conclusion":
            self._render_conclusion_slide(slide, slide_data)
        else:
            self._render_content_slide(slide, slide_data)

    # ─── HELPERS ──────────────────────────────────────────────
    
    def _add_bottom_line(self, slide):
        """Add a horizontal separator line near the bottom (like reference)."""
        line = slide.shapes.add_connector(
            1,  # straight connector
            self.MARGIN_LEFT, Inches(6.14),
            self.MARGIN_LEFT + self.CONTENT_WIDTH, Inches(6.14)
        )
        line.line.color.rgb = self.COLOR_LINE
        line.line.width = Pt(0.75)

    def _add_slide_number(self, slide):
        """Add slide number at bottom-right. Use placeholder if available."""
        slidenum_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type == 13: # SLIDE_NUMBER
                slidenum_ph = ph
                break
        
        if slidenum_ph:
            slidenum_ph.text = str(self.slide_number)
        else:
            txBox = slide.shapes.add_textbox(
                Inches(9.94), Inches(7.23), Inches(3.0), Inches(0.13)
            )
            tf = txBox.text_frame
            tf.text = str(self.slide_number)
            p = tf.paragraphs[0]
            p.font.size = Pt(8)
            p.font.color.rgb = self.COLOR_SUBTITLE
            p.alignment = PP_ALIGN.RIGHT

    def _add_title_bar(self, slide, title_text):
        """Add the section title. Use title placeholder if available, otherwise add textbox."""
        title_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type == 1: # TITLE
                title_ph = ph
                break
        
        if title_ph:
            title_ph.text = title_text
        else:
            # Fallback to looking for PH idx 0 (often title)
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    title_ph = ph
                    break
            
            if title_ph:
                title_ph.text = title_text
            else:
                txBox = slide.shapes.add_textbox(
                    self.MARGIN_LEFT, Inches(0.66),
                    Inches(11.65), Inches(0.58)
                )
                tf = txBox.text_frame
                tf.text = title_text
                p = tf.paragraphs[0]
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = self.COLOR_TITLE

    def _add_subtitle_bar(self, slide, subtitle_text):
        """Add a subtitle below the title."""
        txBox = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(1.24),
            Inches(11.65), Inches(0.3)
        )
        tf = txBox.text_frame
        tf.margin_left = Inches(0.1)
        tf.text = subtitle_text
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = self.COLOR_SUBTITLE

    def _add_numbered_badge(self, slide, number, left, top, size=Inches(0.75)):
        """Add a circular numbered badge (like '01', '02' in reference)."""
        # Background circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, size, size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.COLOR_BADGE_BG
        circle.line.fill.background()
        
        # Number text
        tf = circle.text_frame
        tf.text = f"{number:02d}"
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

    def _add_vertical_divider(self, slide, x, top, height):
        """Add a vertical divider line between columns."""
        line = slide.shapes.add_connector(
            1, x, top, x, top + height
        )
        line.line.color.rgb = self.COLOR_LINE
        line.line.width = Pt(0.5)

    def _add_connecting_arrow(self, slide, x, y, width):
        """Add a horizontal arrow connecting infographic steps."""
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x, y, width, Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.COLOR_LINE
        arrow.line.fill.background()

    def _add_content_card(self, slide, left, top, width, height, title, body, badge_num=None):
        """Add a content card with internal margins and vertical centering."""
        card_top = top
        
        # Add numbered badge if provided
        if badge_num is not None:
            badge_size = Inches(0.55)
            badge_left = left + (width - badge_size) / 2
            self._add_numbered_badge(slide, badge_num, badge_left, card_top, badge_size)
            card_top += Inches(0.7)
        
        # Card title (bold)
        title_box = slide.shapes.add_textbox(left, card_top, width, Inches(0.45))
        tf = title_box.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.LEFT
        
        # Card body text
        body_top = card_top + Inches(0.45)
        body_height = height - (body_top - top)
        body_box = slide.shapes.add_textbox(left, body_top, width, body_height)
        tf2 = body_box.text_frame
        tf2.margin_left = Inches(0.1)
        tf2.margin_right = Inches(0.1)
        tf2.margin_top = Inches(0.05)
        tf2.vertical_anchor = MSO_ANCHOR.TOP
        tf2.word_wrap = True
        tf2.text = body
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(11)
        p2.font.color.rgb = self.COLOR_BODY
        p2.alignment = PP_ALIGN.LEFT

    # ─── SLIDE TYPE RENDERERS ─────────────────────────────────

    def _render_title_slide(self, slide, slide_data: Slide):
        """Cover slide using placeholders 10 and 11."""
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 10:
                ph.text = slide_data.title
                for para in ph.text_frame.paragraphs:
                    para.font.size = Pt(36)
                    para.font.bold = True
            elif idx == 11:
                ph.text = slide_data.subtitle or ""
                for para in ph.text_frame.paragraphs:
                    para.font.size = Pt(16)
        
        # If no placeholders were found, fall back to text boxes
        if not any(True for ph in slide.placeholders):
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
            tf = txBox.text_frame
            tf.text = slide_data.title
            tf.paragraphs[0].font.size = Pt(40)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.COLOR_WHITE
            
            if slide_data.subtitle:
                sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(1))
                stf = sub.text_frame
                stf.text = slide_data.subtitle
                stf.paragraphs[0].font.size = Pt(18)
                stf.paragraphs[0].font.color.rgb = self.COLOR_WHITE

    def _render_content_slide(self, slide, slide_data: Slide):
        """
        Render a content/bullet slide matching reference style:
        - Title bar at top
        - Multi-column grid with numbered badges and dividers
        - Bottom separator line + slide number
        """
        # Title
        self._add_title_bar(slide, slide_data.title)
        
        # Subtitle
        if slide_data.subtitle:
            self._add_subtitle_bar(slide, slide_data.subtitle)
        
        body_items = slide_data.body_groups or []
        
        if len(body_items) == 0:
            # Nothing to render
            pass
        elif len(body_items) <= 2:
            # Two-column layout
            self._render_two_column(slide, body_items)
        elif len(body_items) <= 4:
            # Grid layout with numbered badges (like reference slide 10)
            self._render_grid_cards(slide, body_items)
        else:
            # Bullet list for many items
            self._render_bullet_list(slide, body_items)
        
        # Bottom elements
        self._add_bottom_line(slide)
        self._add_slide_number(slide)

    def _render_two_column(self, slide, items):
        """Render 2 items in a two-column layout with vertical divider."""
        col_width = Inches(5.5)
        top = Inches(1.8)
        height = Inches(4.0)
        
        for i, text in enumerate(items):
            left = self.MARGIN_LEFT + (i * (col_width + Inches(1.5)))
            self._add_content_card(slide, left, top, col_width, height,
                                   text.split('.')[0] if '.' in text else text[:30],
                                   text, badge_num=i+1)
        
        # Vertical divider between columns
        divider_x = self.MARGIN_LEFT + col_width + Inches(0.75)
        self._add_vertical_divider(slide, divider_x, top, height)

    def _render_grid_cards(self, slide, items):
        """Render items in a numbered grid (like reference slide 10 with 01-05 cards)."""
        num_items = len(items)
        gap = Inches(0.3)
        total_gap = gap * (num_items - 1)
        col_width = (self.CONTENT_WIDTH - total_gap) / num_items
        top = Inches(1.8)
        card_height = Inches(4.0)
        
        for i, text in enumerate(items):
            left = self.MARGIN_LEFT + i * (col_width + gap)
            
            # Split text into title and body if possible
            parts = text.split(':', 1) if ':' in text else text.split('.', 1)
            title = parts[0].strip()
            body = parts[1].strip() if len(parts) > 1 else text
            
            self._add_content_card(slide, left, top, col_width, card_height,
                                   title, body, badge_num=i+1)
            
            # Add vertical divider between cards (not after last)
            if i < num_items - 1:
                divider_x = left + col_width + gap / 2
                self._add_vertical_divider(slide, divider_x, Inches(2.0), Inches(4.5))

    def _render_bullet_list(self, slide, items):
        """Render as a styled bullet list with left accent bar and internal margins."""
        # Accent bar on the left
        bar_height = min(Inches(4.2), Inches(0.5 * len(items) + 0.5))
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.MARGIN_LEFT, Inches(1.8),
            Inches(0.08), bar_height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLOR_ACCENT
        bar.line.fill.background()
        
        # Bullet items container
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.5)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)
        tf.word_wrap = True
        
        for i, text in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▸  {text}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.COLOR_BODY
            p.space_after = Pt(12)
            p.level = 0

    def _render_chart_slide(self, slide, slide_data: Slide):
        """Render a chart slide with title, native chart, and caption."""
        # Title
        self._add_title_bar(slide, slide_data.title)
        
        if not slide_data.chart_data:
            self._render_content_slide(slide, slide_data)
            return
            
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = slide_data.chart_data.categories
        
        for series in slide_data.chart_data.series:
            chart_data_obj.add_series(series.name, series.values)
            
        # Map chart type
        c_type = slide_data.chart_data.chart_type
        chart_type_map = {
            "Bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "Pie": XL_CHART_TYPE.PIE,
            "Line": XL_CHART_TYPE.LINE,
        }
        pptx_chart_type = chart_type_map.get(c_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Chart positioning (matching reference slide 9)
        x = Inches(0.62)
        y = Inches(1.66)
        cx = Inches(12.08)
        cy = Inches(5.23)
        
        chart_frame = slide.shapes.add_chart(pptx_chart_type, x, y, cx, cy, chart_data_obj)
        
        # Style the chart
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
        
        # Color each series
        for i, series in enumerate(chart.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = self.CHART_COLORS[i % len(self.CHART_COLORS)]

        # Caption below chart
        if slide_data.subtitle:
            cap_box = slide.shapes.add_textbox(
                Inches(4.36), Inches(7.12), Inches(4.61), Inches(0.25)
            )
            tf = cap_box.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.text = slide_data.subtitle
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.italic = True
            tf.paragraphs[0].font.color.rgb = self.COLOR_SUBTITLE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Bottom elements
        self._add_bottom_line(slide)
        self._add_slide_number(slide)

    def _render_infographic_slide(self, slide, slide_data: Slide):
        """
        Render a process-flow infographic matching reference style:
        - Title at top
        - Numbered columns with chevron shapes
        - Vertical dividers between steps
        - Descriptions below each step
        """
        if not slide_data.process_flow:
            self._render_content_slide(slide, slide_data)
            return
        
        # Title
        self._add_title_bar(slide, slide_data.title)
        
        # Subtitle
        if slide_data.subtitle:
            self._add_subtitle_bar(slide, slide_data.subtitle)

        steps = slide_data.process_flow
        num_steps = len(steps)
        if num_steps == 0:
            return
        
        # Grid layout for the process steps
        gap = Inches(0.3)
        total_gap = gap * (num_steps - 1)
        col_width = (self.CONTENT_WIDTH - total_gap) / num_steps
        top = Inches(1.8)
        
        for i, step in enumerate(steps):
            left = self.MARGIN_LEFT + i * (col_width + gap)
            
            # Numbered badge
            badge_size = Inches(0.75)
            badge_left = left + (col_width - badge_size) / 2
            self._add_numbered_badge(slide, i + 1, badge_left, top, badge_size)
            
            # Chevron shape below badge
            chevron_top = top + Inches(1.0)
            chevron_height = Inches(1.2)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                int(left), int(chevron_top),
                int(col_width), int(chevron_height)
            )
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.CHART_COLORS[i % len(self.CHART_COLORS)]
            shape.line.fill.background()
            
            tf = shape.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.word_wrap = True
            tf.text = step.title
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # Description card below chevron
            desc_top = chevron_top + chevron_height + Inches(0.3)
            desc_box = slide.shapes.add_textbox(
                int(left), int(desc_top),
                int(col_width), Inches(2.0)
            )
            dtf = desc_box.text_frame
            dtf.margin_left = Inches(0.1)
            dtf.margin_right = Inches(0.1)
            dtf.word_wrap = True
            dtf.text = step.description or ""
            dp = dtf.paragraphs[0]
            dp.font.size = Pt(11)
            dp.font.color.rgb = self.COLOR_BODY
            dp.alignment = PP_ALIGN.LEFT
            
            # Vertical divider between steps (not after last)
            if i < num_steps - 1:
                divider_x = left + col_width + gap / 2
                self._add_vertical_divider(slide, divider_x, top, Inches(5.0))
                
                # Add horizontal connecting arrow between chevrons
                arrow_x = left + col_width
                arrow_y = chevron_top + chevron_height / 2 - Inches(0.1)
                self._add_connecting_arrow(slide, arrow_x, arrow_y, gap)
        
        # Bottom elements
        self._add_bottom_line(slide)
        self._add_slide_number(slide)

    def _render_swot_slide(self, slide, slide_data: Slide):
        """Render a 2x2 SWOT analysis grid."""
        self._add_title_bar(slide, slide_data.title)
        
        swot = slide_data.swot_data
        if not swot:
            self._render_content_slide(slide, slide_data)
            return
            
        # Map model attributes to display keys
        swot_map = {
            "STRENGTHS": swot.strengths,
            "WEAKNESSES": swot.weaknesses,
            "OPPORTUNITIES": swot.opportunities,
            "THREATS": swot.threats
        }
        
        card_w = (self.CONTENT_WIDTH - Inches(0.4)) / 2
        card_h = Inches(2.2)
        
        for i, (key, items) in enumerate(swot_map.items()):
            row = i // 2
            col = i % 2
            left = self.MARGIN_LEFT + col * (card_w + Inches(0.4))
            top = self.CONTENT_TOP + row * (card_h + Inches(0.4))
            
            # Card background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, card_w, card_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.COLOR_LIGHT_BG
            shape.line.color.rgb = self.COLOR_ACCENT if row == 0 else self.COLOR_ACCENT2
            shape.line.width = Pt(2)
            
            # Key title
            tBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), card_w - Inches(0.2), Inches(0.4))
            tf = tBox.text_frame
            tf.text = key
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = self.COLOR_TITLE
            
            # Content
            cBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.5), card_w - Inches(0.2), card_h - Inches(0.6))
            ctf = cBox.text_frame
            ctf.word_wrap = True
            
            for item in items:
                p = ctf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(10)
                p.font.color.rgb = self.COLOR_BODY

    def _render_comparison_slide(self, slide, slide_data: Slide):
        """Render a side-by-side comparison grid."""
        self._add_title_bar(slide, slide_data.title)
        
        pairs = slide_data.comparison_data or []
        if not pairs:
            self._render_content_slide(slide, slide_data)
            return
            
        num_cols = len(pairs)
        gap = Inches(0.2)
        col_w = (self.CONTENT_WIDTH - (gap * (num_cols-1))) / num_cols
        
        for i, pair in enumerate(pairs):
            left = self.MARGIN_LEFT + i * (col_w + gap)
            top = self.CONTENT_TOP
            
            # Header
            header_h = Inches(0.6)
            h_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, col_w, header_h
            )
            h_shape.fill.solid()
            h_shape.fill.fore_color.rgb = self.COLOR_ACCENT if i % 2 == 0 else self.COLOR_ACCENT2
            h_shape.line.fill.background()
            
            htf = h_shape.text_frame
            htf.text = pair.key
            htf.vertical_anchor = MSO_ANCHOR.MIDDLE
            htf.paragraphs[0].alignment = PP_ALIGN.CENTER
            htf.paragraphs[0].font.bold = True
            htf.paragraphs[0].font.color.rgb = self.COLOR_WHITE
            
            # Content
            content_h = Inches(3.5)
            c_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top + header_h, col_w, content_h
            )
            c_shape.fill.solid()
            c_shape.fill.fore_color.rgb = self.COLOR_LIGHT_BG
            c_shape.line.color.rgb = self.COLOR_LINE
            
            ctf = c_shape.text_frame
            ctf.margin_left = Inches(0.1)
            ctf.margin_right = Inches(0.1)
            ctf.text = pair.value
            ctf.paragraphs[0].font.size = Pt(11)
            ctf.paragraphs[0].font.color.rgb = self.COLOR_BODY
            
        self._add_bottom_line(slide)
        self._add_slide_number(slide)

    def _render_conclusion_slide(self, slide, slide_data: Slide):
        """Render conclusion/thank you slide. The Thank You layout is a pre-designed 
        background with no placeholders - we just use it as-is."""
        # The Thank You layout has 0 placeholders and already contains
        # the visual "Thank You" design in its background. Adding text
        # on top causes jumbled/overlapping content. So we leave it clean.
        pass
