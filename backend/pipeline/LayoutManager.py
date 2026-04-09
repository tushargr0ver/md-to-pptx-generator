from logging import getLogger
from pptx import Presentation
from lxml import etree

logger = getLogger(__name__)

class LayoutManager:
    def __init__(self, master_path: str):
        self.prs = Presentation(master_path)
        self._remove_existing_slides()
        self._detect_layouts()
        
    def _remove_existing_slides(self):
        """Remove ALL pre-existing slides from the template so we start clean."""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            )
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])

    def _detect_layouts(self):
        """Auto-detect and log all layouts and their placeholders."""
        self.layout_info = {}
        for i, layout in enumerate(self.prs.slide_layouts):
            placeholders = {}
            for ph in layout.placeholders:
                placeholders[ph.placeholder_format.idx] = {
                    'name': ph.name,
                    'type': str(ph.placeholder_format.type),
                }
            self.layout_info[i] = {
                'name': layout.name,
                'placeholders': placeholders,
            }
            logger.info(f"Layout {i}: '{layout.name}' -> {placeholders}")

        # Map our slide types to the best available layouts
        # Based on the Accenture template structure:
        # 0: "1_Cover"      - Title/cover slide (placeholders 10, 11)
        # 1: "2_Cover"      - Alternative cover (placeholders 10, 11)  
        # 2: "Divider"      - Section divider (placeholder 10)
        # 3: "Blank"        - Completely blank
        # 4: "Title only"   - Title + content area (placeholders 0, 11, 4)
        # 5: "Thank You"    - Closing slide
        self.layout_map = {
            "title_slide": 0,            # Cover slide
            "section_divider": 2,        # Divider
            "content_text": 3,           # Blank - full programmatic control
            "bullet_points": 3,          # Blank - full programmatic control
            "content_chart": 3,          # Blank - full programmatic control
            "infographic_process": 3,    # Blank - full programmatic control
            "conclusion": 5,             # Thank You
        }
        
    def get_layout(self, slide_type: str):
        """Returns the appropriate slide layout for a given slide_type."""
        layout_index = self.layout_map.get(slide_type, 4)  # Default to Title Only
        try:
            return self.prs.slide_layouts[layout_index]
        except IndexError:
            return self.prs.slide_layouts[0]
    
    def get_placeholder_indices(self, slide_type: str):
        """Returns the placeholder indices for a given slide type."""
        layout_index = self.layout_map.get(slide_type, 4)
        info = self.layout_info.get(layout_index, {})
        return info.get('placeholders', {})
